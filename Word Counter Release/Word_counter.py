from collections import defaultdict
from datetime import datetime
import glob
import os
from turtle import title
import pandas as pd
import aspose.slides as slides
import aspose.pydrawing as drawing
from styleframe import StyleFrame
from pptx import Presentation
import docx
from PyPDF2 import PdfFileReader
from kivy.app import App
from kivy.uix.button import Button
from kivy.uix.label import Label
from kivy.uix.boxlayout import BoxLayout
from kivy.core.window import Window
from kivy.uix.gridlayout import GridLayout
from kivy.uix.popup import Popup

Window.size = (800, 200)
Window.clearcolor = (255 / 255, 165 / 255, 3 / 255, 1)
Window.title = "Word Counter"


class WordCounterApp(App):
    def calculate_data_of_running(self):
        time = datetime.now()
        return time

    def check_existing_file(self, file_name):
        if not os.path.exists(file_name):
            print("File does not exsist " + file_name)

    def calculate_amount_of_running(self, file_name):
        self.check_existing_file(file_name)
        try:
            global amount_of_running
            c = open(file_name, "r+")
            amount_of_running = int(c.readline())
            c.seek(0)
            c.truncate()
            amount_of_running += 1
            c.write(str(amount_of_running))
            c.close()
        except FileNotFoundError:
            print("File does not open " + file_name)

    def read_txt_files(self, folder, file_list):
        for filename in glob.glob(folder + "*.txt"):
            with open(os.path.join(os.getcwd(), filename), "r") as f:

                words_list = f.read().split()

                for word in words_list:
                    file_list.append(word)
        return file_list

    def read_pptx_files(self, folder, file_list):
        for filename in glob.glob(folder + "*.pptx"):
            ppt = Presentation(filename)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text
                            words_list = text.split()
                            for word in words_list:
                                file_list.append(word)
        return file_list

    def read_docx_files(self, folder, file_list):
        for filename in glob.glob(folder + "*.docx"):
            file = docx.Document(filename)
            for paragraph in file.paragraphs:
                text = paragraph.text
                words_list = text.split()
                for word in words_list:
                    file_list.append(word)
        return file_list

    def read_pdf_files(self, folder, file_list):
        global_text = []
        for filename in glob.glob(folder + "*.pdf"):
            with open(filename, "rb") as filehandle:
                pdf = PdfFileReader(filehandle)
                pages = pdf.getNumPages()
                for i in range(pages):
                    page = pdf.getPage(i)
                    global_text.append(page.extractText())
                words_list = global_text.split()
                for word in words_list:
                    file_list.append(word)
        return file_list

    def read_all_files(self, folder, file_list):
        folder = folder + "\\source files\\"
        file_list = self.read_txt_files(folder, file_list)
        file_list = self.read_pptx_files(folder, file_list)
        file_list = self.read_docx_files(folder, file_list)
        file_list = self.read_pdf_files(folder, file_list)
        return file_list

    def to_lower_register(self, words_list):
        for word in range(len(words_list)):
            words_list[word] = words_list[word].lower()
        return words_list

    def delete_outsiders_symbols(self, words_list, symbols_list):
        for word in range(len(words_list)):
            for symbol in range(len(symbols_list)):
                words_list[word] = words_list[word].replace(symbols_list[symbol], "")
        return words_list

    def calculate_amount_of_words(self, words_list):
        word_count = defaultdict(int)
        for word in words_list:
            if word in word_count:
                word_count[word] += 1
            else:
                word_count[word] = 1
        return word_count

    def write_result(self, name_result_file, start_data, dictionary):
        self.check_existing_file(name_result_file)
        try:
            with open(name_result_file, "w+", encoding="utf-8") as out:
                out.write(
                    "Data of start "
                    + str(start_data)
                    + " "
                    + "Amount of running programm "
                    + str(amount_of_running)
                    + "\n"
                )
                for key, val in dictionary.items():
                    out.write("{}:{}\n".format(key, val))
        except FileNotFoundError:
            print("File does not exsist " + name_result_file)

    def sort_dictionary(self, dictionary):
        sorted_dict = {}
        sorted_keys = sorted(dictionary, key=dictionary.get, reverse=True)

        for key in sorted_keys:
            sorted_dict[key] = dictionary[key]

        return sorted_dict

    def to_exl(self, folder, start_data, dictionary):
        df = pd.DataFrame(
            {
                "Amount of running programm ": amount_of_running,
                "Data of start ": str(start_data),
                "Words ": dictionary.keys(),
                "Number of repetitions ": dictionary.values(),
            }
        )
        df.to_excel(folder + "\\result.xlsx", sheet_name="Amount of word", index=False)

    def main_word_counter(self):
        amount_of_running = 0

        words_list_from_files = []
        symbols_list = [",", ".", "(", ")", ";", ":", "!", "*", "␋", "␋"]

        current_folder = os.path.abspath(__file__).rpartition("\\")[0]
        time = self.calculate_data_of_running()
        self.calculate_amount_of_running(current_folder + "\\count.txt")
        words_list_from_files = self.read_all_files(
            current_folder, words_list_from_files
        )
        words_list_from_files = self.to_lower_register(words_list_from_files)
        words_list_from_files = self.delete_outsiders_symbols(
            words_list_from_files, symbols_list
        )
        word_count = self.calculate_amount_of_words(words_list_from_files)
        word_count = self.sort_dictionary(word_count)
        self.write_result(current_folder + "\\result.txt", time, word_count)
        self.to_exl(current_folder, time, word_count)
        print("Programm finished work")

    def btn_run_pressed(self, instance):
        print("Button is pressed")
        self.main_word_counter()
        label = Label(text="Word counter completed successfully!")
        popupWindow = Popup(
            title="Message", content=label, size_hint=(None, None), size=(350, 100)
        )
        popupWindow.open()

    def build(self):
        layout = GridLayout(cols=2)
        box = BoxLayout()
        btn_run = Button(text="Run")
        btn_run.bind(on_press=self.btn_run_pressed)
        label = Label(
            text='Hello, this is a file counting program!\nThe program can count words in such formats as: .pdf, .docx, .pptx, .txt for this:\n1) Put the files in a folder\n2)Press the "Run" button\n3) Look at the result in the file "result.x" that appears in the folder with the program',
            size_hint_x=None,
            font_size=15,
            width=650,
            color=(0 / 255, 0 / 255, 0 / 255),
        )
        box.add_widget(label)
        box.add_widget(btn_run)

        return box


if __name__ == "__main__":
    app = WordCounterApp()
    app.run()
